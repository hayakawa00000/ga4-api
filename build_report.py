"""
build_report.py - GA4/GSC/CV レポート PowerPoint生成スクリプト
テンプレート: 26年〇月デジマアクセル月次レポート_〇〇（〇〇店）.pptx
  Slide1: 表紙（店舗名・期間を書き換え）
  Slide2: 目次（固定）
  Slide3: 本文スライド雛形 → P3以降はこれをコピーして使用

usage: python build_report.py report_data_test.json [output.pptx]
"""
import json
import sys
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_TICK_LABEL_POSITION
from pptx.oxml.ns import qn

import plot_utils

# ============================================================
# テンプレートファイル
# ============================================================
TEMPLATE_FILE = "template.pptx"

# ============================================================
# 定数（7.5" x 5.62" = 4:3 スライド）
# ============================================================
SW = Inches(7.5)   # Slide Width
SH = Inches(5.62)  # Slide Height

# カラー
C_BG_HEADER = RGBColor(0x26, 0x41, 0x6E)   # ヘッダー紺（テンプレートのアクセント色）
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT      = RGBColor(0x1A, 0x1A, 0x2E)
C_SUBTEXT   = RGBColor(0x55, 0x55, 0x66)
C_GRAY      = RGBColor(0xBB, 0xBB, 0xBB)
C_LIGHT     = RGBColor(0xF0, 0xF4, 0xFA)
C_PRIMARY   = RGBColor(0x26, 0x41, 0x6E)
C_ACCENT    = RGBColor(0x2E, 0x86, 0xAB)
C_RED       = RGBColor(0xE5, 0x53, 0x3C)
C_ORANGE    = RGBColor(0xF5, 0xA5, 0x23)
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
C_BLUE      = RGBColor(0x2E, 0x86, 0xAB)

BADGE_COLORS = {
    "増加傾向": (RGBColor(0xE8, 0xF5, 0xE9), C_GREEN),
    "大幅改善": (RGBColor(0xE3, 0xF2, 0xFD), C_BLUE),
    "要改善":   (RGBColor(0xFF, 0xEB, 0xEE), C_RED),
    "目標達成": (RGBColor(0xE8, 0xF5, 0xE9), C_GREEN),
    "回復傾向": (RGBColor(0xE8, 0xF5, 0xE9), C_GREEN),
    "横ばい":   (RGBColor(0xFA, 0xFA, 0xFA), C_GRAY),
    "悪化":     (RGBColor(0xFF, 0xEB, 0xEE), C_RED),
    "未達":     (RGBColor(0xFF, 0xEB, 0xEE), C_RED),
}

PROPOSAL_COLORS = {
    "red":    (RGBColor(0xFF, 0xEB, 0xEE), C_RED,    "🔴 最優先"),
    "orange": (RGBColor(0xFF, 0xF3, 0xE0), C_ORANGE, "🟠 優先"),
    "blue":   (RGBColor(0xE3, 0xF2, 0xFD), C_BLUE,   "🔵 適宜"),
}

FONT_JP = "Meiryo UI"

# ============================================================
# ユーティリティ
# ============================================================
def rgb(r, g, b): return RGBColor(r, g, b)
def inch(v): return Inches(v)

def set_text(shape, text, font_size=None, bold=None, color=None,
             align=PP_ALIGN.LEFT, italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # 既存runs削除
    for run in p.runs:
        run._r.getparent().remove(run._r)
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_JP
    if font_size:  run.font.size  = Pt(font_size)
    if bold is not None: run.font.bold = bold
    if italic: run.font.italic = True
    if color:  run.font.color.rgb = color

def add_text(slide, left, top, width, height, text="",
             font_size=10, bold=False, color=None, bg=None,
             align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    tb = slide.shapes.add_textbox(inch(left), inch(top), inch(width), inch(height))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
    return tb

def add_rect(slide, left, top, width, height, fill, line_color=None, lw=0.5):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, inch(left), inch(top), inch(width), inch(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def add_table(slide, left, top, width, height, headers, rows,
              hdr_bg=None, fs=8, col_widths=None):
    if not rows: return
    hdr_bg = hdr_bg or C_PRIMARY
    nc = len(headers); nr = len(rows) + 1
    tbl = slide.shapes.add_table(nr, nc, inch(left), inch(top), inch(width), inch(height)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = inch(w)
    # ヘッダー
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = hdr_bg
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True; run.font.size = Pt(fs)
                run.font.color.rgb = C_WHITE; run.font.name = FONT_JP
    # データ
    for i, row in enumerate(rows):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val) if val is not None else "-"
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(fs); run.font.name = FONT_JP
                    run.font.color.rgb = C_TEXT
    return tbl

def delta_str(v):
    if v is None: return "-"
    return f"+{v:,}" if v > 0 else f"{v:,}"

# ============================================================
# スライドコピー（テンプレートSlide3をベースに新スライド追加）
# ============================================================
def add_content_slide(prs, template_slide):
    """テンプレートのSlide3をdeepcopyして新スライドを追加"""
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # テンプレートのshapes XMLを複製して追加
    sp_tree = new_slide.shapes._spTree
    # 既存のspTreeをクリア（タイトルプレースホルダ含む）
    for el in list(sp_tree):
        sp_tree.remove(el)
    # テンプレートのspTreeをコピー
    for el in template_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))
    return new_slide

C_TITLE_RED = RGBColor(0xC0, 0x30, 0x20)  # テンプレートのタイトル赤色

def slide_header(slide, title, period=""):
    """テンプレートと同じヘッダー: タイトルテキストを赤色に"""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "タイトル 1":
            set_text(shape, title, font_size=14, bold=True, color=C_TITLE_RED)
            return
    add_text(slide, 0.17, 0.08, 5.5, 0.26, title,
             font_size=14, bold=True, color=C_TITLE_RED)

# ============================================================
# データロード・summary数値補完
# ============================================================
def load_data(path):
    with open(path, encoding="utf-8") as f:
        d = json.load(f)
    ga = d.get("ga4_monthly", [])
    cv = d.get("cv_months", [])
    s  = d.get("summary", {})
    if len(ga) >= 2:
        cur = ga[-1]; prv = ga[-2]
        ss = s.setdefault("sessions", {})
        ss.setdefault("value",    cur["sessions"])
        ss.setdefault("prev",     prv["sessions"])
        ss.setdefault("diff_pct", round((cur["sessions"]-prv["sessions"])/max(prv["sessions"],1)*100,1))
        ar = s.setdefault("area_sessions", {})
        ar.setdefault("value",     cur["area_sessions"])
        ar.setdefault("prev",      prv["area_sessions"])
        ar.setdefault("diff_pct",  round((cur["area_sessions"]-prv["area_sessions"])/max(prv["area_sessions"],1)*100,1))
        ar.setdefault("rate",      cur["area_rate"])
        ar.setdefault("prev_rate", prv["area_rate"])
        iq = s.setdefault("inquiry", {})
        iq.setdefault("value",     cur["inquiry_views"])
        iq.setdefault("prev",      prv["inquiry_views"])
        iq.setdefault("diff_pct",  round((cur["inquiry_views"]-prv["inquiry_views"])/max(prv["inquiry_views"],1)*100,1))
        iq.setdefault("rate",      cur["inquiry_rate"])
        iq.setdefault("prev_rate", prv["inquiry_rate"])
    if len(cv) >= 2:
        cc = s.setdefault("cv", {})
        cc.setdefault("value", cv[-1]["actual"])
        cc.setdefault("prev",  cv[-2]["actual"])
        cc.setdefault("diff",  cv[-1]["actual"] - cv[-2]["actual"])
    d["summary"] = s
    return d

# ============================================================
# P1: 表紙（Slide1を編集）
# ============================================================
def edit_p1_cover(slide, d):
    """P1表紙: テキストのみ差し替え（フォント・色・サイズは一切変更しない）"""
    store = d.get("store_name", "○○様")
    period = d.get("period", "○年○月")
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        txt = shape.text_frame.text
        if "〇〇様" in txt or "○○様" in txt:
            # runのtextのみ変更、スタイルは保持
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "〇〇様" in run.text or "○○様" in run.text:
                        run.text = f"{store}様"
        elif "プロタイムズ〇〇店" in txt:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "〇〇店" in run.text:
                        run.text = f"（プロタイムズ {store}）"
        elif "20〇年〇月" in txt or "アクセル定例" in txt:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "〇" in run.text or "アクセル定例" in run.text:
                        run.text = f"{period}　アクセル定例MTG"

# ============================================================
# P3: CV・CPA実績比較（参考レポートP1の座標系を再現）
# ============================================================
def build_p3_cv(slide, d):
    slide_header(slide, "実績比較（CV・CPA）")
    period_1st = d.get("period_1st", "")
    if period_1st:
        import re
        m = re.search(r'(\d+)年(\d+)月.*?(\d+)-(\d+)-(\d+)', period_1st)
        if m:
            p_str = f"{m.group(1)}/{m.group(2)}/1～{m.group(3)}/{int(m.group(4))}/{int(m.group(5))}"
        else:
            p_str = period_1st.replace(" - ", "～").replace("-", "～")
        add_text(slide, 3.2, 0.10, 4.0, 0.25, f"【対象期間：{p_str}】",
                 font_size=11, bold=True, color=C_TEXT)
    cv = d.get("cv_months", [])

    cv_rev = list(reversed(cv))
    
    # --- 3ヶ月カード（参考座標: W=2.15 H=1.88 T=0.76 スライド幅7.5"以内）---
    card_lefts = [0.35, 2.68, 5.0]
    # 新しい月(当月)を左にするため色も反転：緑 -> 青 -> 薄青
    card_colors = [rgb(0xE8,0xF5,0xE9), rgb(0xE3,0xF2,0xFD), rgb(0xEE,0xF2,0xFB)]
    for i, m in enumerate(cv_rev):
        lx = card_lefts[i]; ty = 0.76
        add_rect(slide, lx, ty, 2.1, 1.88, card_colors[i])
        add_text(slide, lx+0.05, ty+0.06, 1.95, 0.22, m.get("ym",""),
                 font_size=10, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
        actual = m.get("actual",0); target = m.get("target",0)
        result = "✅ 達成" if actual >= target else "❌ 未達"
        rc = C_GREEN if actual >= target else C_RED
        add_text(slide, lx+0.05, ty+0.3, 1.95, 0.24, result,
                 font_size=11, bold=True, color=rc, align=PP_ALIGN.CENTER)
        items = [("目標CV", f"{target}件"), ("実績CV", f"{actual}件"),
                 ("広告費", f"¥{m.get('budget',0):,}"), ("CPA", f"¥{m.get('cpa',0):,}")]
        for j,(k,v) in enumerate(items):
            yt = ty + 0.6 + j*0.3
            cval = C_TEXT 
            if k == "CPA":
                cpa_val = m.get('cpa',0)
                cval = C_GREEN if cpa_val <= 50000 else C_RED
            add_text(slide, lx+0.08, yt, 1.0, 0.25, k+":", font_size=9, color=C_SUBTEXT)
            add_text(slide, lx+1.05, yt, 0.9, 0.25, v, font_size=9, bold=True,
                     color=cval, align=PP_ALIGN.RIGHT)

    # --- CVテーブル（下部）スライド幅7.15"以内 ---
    headers = ["月", "目標CV", "実績CV", "広告費", "CPA"]
    rows = [[m["ym_short"], m["target"], m["actual"],
             f"¥{m['budget']:,}", f"¥{m['cpa']:,}"] for m in cv_rev]
    add_table(slide, 0.2, 2.74, 7.1, 1.1, headers, rows, fs=9,
              col_widths=[1.1, 1.2, 1.2, 2.4, 1.2])

    # --- コメント（ロゴ回避：Hを小さく） ---
    cv_comment = d.get("cv_comment", "")
    add_rect(slide, 0.35, 4.0, 6.8, 1.3, C_LIGHT)
    add_text(slide, 0.45, 4.05, 1.0, 0.25, "📝 コメント",
             font_size=9, bold=True, color=C_PRIMARY)
    add_text(slide, 0.45, 4.3, 6.5, 0.95, cv_comment,
             font_size=10, color=C_TEXT)


# ============================================================
# P4: 全体サマリー（参考レポートP2の座標系を再現）
# ============================================================
def build_p4_summary(slide, d):
    slide_header(slide, "全体サマリー")
    period_1st = d.get("period_1st", "")
    if period_1st:
        import re
        m = re.search(r'(\d+)年(\d+)月.*?(\d+)-(\d+)-(\d+)', period_1st)
        if m:
            p_str = f"{m.group(1)}/{m.group(2)}/1～{m.group(3)}/{int(m.group(4))}/{int(m.group(5))}"
        else:
            p_str = period_1st.replace(" - ", "～").replace("-", "～")
        add_text(slide, 1.8, 0.10, 4.0, 0.25, f"【対象期間：{p_str}】",
                 font_size=11, bold=True, color=C_TEXT)
    s = d.get("summary", {})

    def card(lx, ty, title, value_str, sub_str, pct_str, badge, badge_sub):
        bg, fg = BADGE_COLORS.get(badge, (C_LIGHT, C_TEXT))
        add_rect(slide, lx, ty, 3.31, 1.71, C_WHITE,
                 line_color=rgb(0xDD,0xDD,0xEE), lw=0.5)
        
        # Title
        add_text(slide, lx+0.1, ty+0.05, 3.0, 0.22, title, font_size=9, color=C_PRIMARY)
        
        # Value
        add_text(slide, lx+0.1, ty+0.28, 1.5, 0.38, str(value_str), font_size=20, bold=True, color=C_PRIMARY)
        
        # Sub string (e.g. rate or prev month)
        add_text(slide, lx+0.1, ty+0.7, 1.5, 0.2, sub_str, font_size=9, color=C_SUBTEXT)
        
        # Pct string (e.g. 前月比 -70.0%) - put it on the right
        add_text(slide, lx+1.7, ty+0.25, 1.5, 0.22, pct_str, font_size=10, bold=True, color=C_PRIMARY)
        
        # Badge
        add_rect(slide, lx+1.7, ty+0.5, 0.9, 0.28, bg)
        add_text(slide, lx+1.7, ty+0.5, 0.9, 0.28, badge, font_size=9, bold=True, color=fg, align=PP_ALIGN.CENTER)
        
        # Long sub text spanning at bottom
        add_text(slide, lx+0.1, ty+1.0, 3.1, 0.6, badge_sub, font_size=9, color=C_SUBTEXT)

    # 4カード: 左上/右上/左下/右下
    ss = s.get("sessions", {})
    card(0.35, 0.76, "セッション数",
         f"{ss.get('value', '-'):,}" if isinstance(ss.get('value'),int) else "-",
         f"前月: {ss.get('prev','-'):,}" if isinstance(ss.get('prev'),int) else "",
         f"前月比 {ss.get('diff_pct',0):+.1f}%",
         ss.get("badge",""), ss.get("badge_sub",""))

    ar = s.get("area_sessions", {})
    card(3.84, 0.76, "商圏内アクセス数",
         f"{ar.get('value','-'):,}" if isinstance(ar.get('value'),int) else "-",
         f"率: {ar.get('prev_rate','-')}%→{ar.get('rate','-')}%",
         f"前月比 {ar.get('diff_pct',0):+.1f}%",
         ar.get("badge",""), ar.get("badge_sub",""))

    iq = s.get("inquiry", {})
    card(0.35, 2.55, "問合せページ遷移数",
         f"{iq.get('value','-'):,}" if isinstance(iq.get('value'),int) else "-",
         f"率: {iq.get('prev_rate','-')}%→{iq.get('rate','-')}%",
         f"前月比 {iq.get('diff_pct',0):+.1f}%",
         iq.get("badge",""), iq.get("badge_sub",""))

    cv_s = s.get("cv", {})
    card(3.84, 2.55, "実CV数",
         f"{cv_s.get('value','-')}件" if isinstance(cv_s.get('value'),int) else "-",
         f"前月: {cv_s.get('prev','-')}件",
         f"前月比 {'+' if (cv_s.get('diff',0) or 0)>=0 else ''}{cv_s.get('diff',0)}件",
         cv_s.get("badge",""), cv_s.get("badge_sub",""))

    # サマリー帯（ロゴ回避: Hを小さく）
    analysis_text = d.get("analysis",{}).get("analysis_text","")
    add_rect(slide, 0.35, 4.35, 6.8, 1.0, C_LIGHT)
    add_text(slide, 0.55, 4.38, 1.2, 0.24, "📊 サマリー分析",
             font_size=9, bold=True, color=C_PRIMARY)
    add_text(slide, 0.55, 4.60, 6.4, 0.70, analysis_text,
             font_size=10, color=C_TEXT)


# ============================================================
# P5: 詳細サマリー（全体指標 / 流入元5分類 / 商圏・問合・デバイス 各3ヶ月 新しい月を上に）
# ============================================================
def build_p5_detail(slide, d):
    slide_header(slide, "詳細サマリー")
    ga = d.get("ga4_monthly", [])
    ga_rev = list(reversed(ga))  # 新しい月を上に

    # テーブル1: 全体指標
    add_text(slide, 0.35, 0.50, 4.0, 0.25, "■ 全体指標（3ヶ月）",
             font_size=9, bold=True, color=C_PRIMARY)
    h1 = ["月","セッション","ユーザー","PV","直帰率","滞在時間"]
    r1 = [[m["ym_short"], f"{m['sessions']:,}", f"{m['users']:,}",
           f"{m['pvs']:,}",
           f"{m['bounce']}%", m["duration"]] for m in ga_rev]
    add_table(slide, 0.35, 0.72, 6.9, 0.9, h1, r1, fs=9,
              col_widths=[0.85,1.15,1.15,1.15,1.2,1.4])

    # テーブル2: 流入元5分類
    add_text(slide, 0.35, 1.72, 4.0, 0.25, "■ 流入元内訳（5分類）",
             font_size=9, bold=True, color=C_PRIMARY)
    h2 = ["月","自然流入","広告流入","直接流入","被リンク流入","SNS流入"]
    r2 = [[m["ym_short"], m["organic"], m["cpc"],
           m["direct"], m["referral"], m["social"]] for m in ga_rev]
    add_table(slide, 0.35, 1.94, 6.7, 0.9, h2, r2, fs=9,
              col_widths=[0.85,1.1,1.1,1.1,1.1,1.15])

    # テーブル3: 商圏・問合・デバイス
    add_text(slide, 0.35, 2.94, 4.0, 0.25, "■ 商圏・問い合わせ・デバイス",
             font_size=9, bold=True, color=C_PRIMARY)
    h3 = ["月","商圏セッション","商圏率","問合遷移","問合率","モバイル%","PC%"]
    r3 = [[m["ym_short"], f"{m['area_sessions']:,}", f"{m['area_rate']}%",
           f"{m['inquiry_views']:,}", f"{m['inquiry_rate']}%",
           f"{m.get('mobile','-')}%", f"{m.get('desktop','-')}%"] for m in ga_rev]
    add_table(slide, 0.35, 3.16, 6.9, 0.9, h3, r3, fs=9,
              col_widths=[0.85,1.05,0.75,1.0,0.75,1.1,1.0])

# ============================================================
# P6: 月別GA4分析（問合せページ内訳 / 商圏市区町村別SS / キーイベント 各3ヶ月結合テーブル）
# ============================================================
def build_p6_ga4(slide, d):
    slide_header(slide, "月別GA4分析")
    ga = d.get("ga4_monthly",[])
    ga_rev = list(reversed(ga))  # 新しい月を左/上に

    # --- 問い合わせページ内訳（結合テーブル）---
    add_text(slide, 0.35, 0.50, 3.3, 0.25, "■ 問い合わせページ内訳",
             font_size=9, bold=True, color=C_PRIMARY)
    all_paths = []
    seen_paths = set()
    for m in ga:
        for it in m.get("inquiry_details", []):
            if it["path"] not in seen_paths:
                all_paths.append(it["path"]); seen_paths.add(it["path"])
    h_iq = ["パス"] + [m["ym_short"] for m in ga_rev]
    r_iq = []
    for path in all_paths:
        row = [path]
        for m in ga_rev:
            pm = {it["path"]: it["views"] for it in m.get("inquiry_details", [])}
            v = pm.get(path, 0)
            row.append(f"{v:,}" if v else "-")
        r_iq.append(row)
    add_table(slide, 0.35, 0.75, 3.2, 0.88, h_iq, r_iq[:3], fs=9,
              col_widths=[1.55]+[0.55]*len(ga))

    # --- 商圏内 市区町村別SS（列順 = 新しい月が左）---
    add_text(slide, 0.35, 1.97, 3.3, 0.25, "■ 商圏内 市区町村別SS",
             font_size=9, bold=True, color=C_PRIMARY)
    all_cities = []
    seen = set()
    for m in ga:
        for c in m.get("area_by_city",[]):
            if c["city"] not in seen:
                all_cities.append(c["city"]); seen.add(c["city"])
    h_city = ["市区町村"] + [m["ym_short"] for m in ga_rev]
    r_city = []
    for city in all_cities:
        row = [city]
        for m in ga_rev:
            cm = {c["city"]: c["sessions"] for c in m.get("area_by_city",[])}
            row.append(f"{cm.get(city,0):,}")
        r_city.append(row)
    add_table(slide, 0.35, 2.20, 3.2, 1.0, h_city, r_city[:8], fs=9,
              col_widths=[1.55]+[0.55]*len(ga))

    # --- キーイベント内訳（結合テーブル）---
    add_text(slide, 3.8, 0.50, 3.3, 0.25, "■ キーイベント内訳",
             font_size=9, bold=True, color=C_PRIMARY)
    all_events = []
    seen_events = set()
    for m in ga:
        for it in m.get("key_events", {}).get("details", []):
            if it["event_name"] not in seen_events:
                all_events.append(it["event_name"]); seen_events.add(it["event_name"])
    h_ke = ["イベント名"] + [m["ym_short"] for m in ga_rev]
    r_ke = []
    # 合計行
    row_total = ["【合計】"]
    for m in ga_rev:
        t = m.get("key_events", {}).get("total", 0)
        row_total.append(f"{t:,}")
    r_ke.append(row_total)
    # 各イベント行
    for ev_name in all_events:
        row = [ev_name]
        for m in ga_rev:
            em = {it["event_name"]: it["count"] for it in m.get("key_events", {}).get("details", [])}
            v = em.get(ev_name, 0)
            row.append(f"{v:,}" if v else "-")
        r_ke.append(row)
    
    add_table(slide, 3.8, 0.75, 3.2, 1.4, h_ke, r_ke[:15], fs=9,
              col_widths=[1.55]+[0.55]*len(ga))


# ============================================================
# P7: 月別GSC分析
# ============================================================
def build_p7_gsc(slide, d):
    slide_header(slide, "月別サーチコンソール分析")
    gsc = d.get("gsc_monthly",[])
    gsc_rev = list(reversed(gsc))  # 新しい月を上に

    add_text(slide, 0.35, 0.50, 4.0, 0.25, "■ 全体指標",
             font_size=9, bold=True, color=C_PRIMARY)
    h1 = ["月","クリック","表示回数","CTR(%)","平均順位"]
    r1 = [[m["ym_short"], f"{m['clicks']:,}", f"{m['impressions']:,}",
           f"{m['ctr']:.2f}%", f"{m['position']:.1f}"] for m in gsc_rev]
    add_table(slide, 0.35, 0.72, 6.8, 0.88, h1, r1, fs=9,
              col_widths=[0.85,1.4,1.7,1.15,1.5])

    if len(gsc) >= 2:
        for idx, (m, label) in enumerate([(gsc[-1],"当月"), (gsc[-2],"前月")]):
            lx = 0.35 + idx * 3.55
            add_text(slide, lx, 1.72, 3.4, 0.25,
                     f"■ クエリ Top（{m['ym_short']}:{label}）",
                     font_size=9, bold=True, color=C_PRIMARY)
            h_q = ["クエリ","CL","Imp","CTR","順位"]
            r_q = [[q["query"], q["clicks"], q["imps"],
                    f"{q['ctr']:.1f}%", f"{q['pos']:.1f}"]
                   for q in m.get("queries",[])]
            add_table(slide, lx, 1.97, 3.4, 3.35, h_q, r_q[:7], fs=9,
                      col_widths=[1.3,0.45,0.55,0.6,0.5])

# ============================================================
# P7.5: 商圏サーチコンソール分析
# ============================================================
def build_p7_5_gsc_area(slide, d, new_slide_fn=None):
    slide_header(slide, "商圏サーチコンソール分析")
    area_gsc = d.get("gsc_area_monthly", [])
    if not area_gsc:
        add_text(slide, 0.35, 1.0, 5.0, 0.5, "データがありません", font_size=12, color=C_TEXT)
        return

    # GSC Area Date format assumption: newest first, index 0 is 当月, index 1 is 前月
    valid_months = [(area_gsc[0], "当月")]
    if len(area_gsc) > 1:
        valid_months.append((area_gsc[1], "前月"))

    all_rows = []
    for m, label in valid_months:
        r_q = []
        for a in m.get("areas", []):
            for q in a.get("queries", []):
                r_q.append([
                    q["query"], 
                    str(q["clicks"]), 
                    str(q["impressions"]), 
                    f"{q['ctr']:.1f}%", 
                    f"{q['position']:.1f}"
                ])
        all_rows.append((m, label, r_q))
        
    max_rows = max([len(rows) for _, _, rows in all_rows]) if all_rows else 0
    pages = (max_rows + 7) // 8
    if pages == 0: pages = 1

    current_slide = slide
    for page_idx in range(pages):
        if page_idx > 0:
            if new_slide_fn:
                current_slide = new_slide_fn()
                slide_header(current_slide, "商圏サーチコンソール分析（続き）")
            else:
                break
                
        start_idx = page_idx * 8
        end_idx = start_idx + 8
        
        for idx, (m, label, r_q) in enumerate(all_rows):
            lx = 0.35 + idx * 3.55
            add_text(current_slide, lx, 0.50, 3.4, 0.25,
                     f"■ 商圏クエリ実績（{m.get('ym_short','')}:{label}）",
                     font_size=9, bold=True, color=C_PRIMARY)
            
            h_q = ["クエリ","CL","Imp","CTR","順位"]
            page_rows = r_q[start_idx:end_idx]
                    
            if page_rows:
                add_table(current_slide, lx, 0.72, 3.4, 4.3, h_q, page_rows, fs=9,
                          col_widths=[1.3,0.45,0.55,0.6,0.5])
            else:
                add_text(current_slide, lx, 0.72, 3.4, 0.5, "表示するクエリがありません", font_size=9, color=C_TEXT)

# ============================================================
# P8: 当月分析
# ============================================================
def build_p8_analysis(slide, d):
    slide_header(slide, "当月分析")
    analysis = d.get("analysis",{})

    secs = [
        ("✅ 良かった点",   analysis.get("good",[]),   C_GREEN, rgb(0xE8,0xF5,0xE9)),
        ("⚠️ 課題・改善点", analysis.get("issues",[]), C_RED,   rgb(0xFF,0xEB,0xEE)),
    ]
    for idx, (title, items, fg, bg) in enumerate(secs):
        lx = 0.2 + idx * 3.7
        add_rect(slide, lx, 0.62, 3.5, 4.8, bg)
        add_text(slide, lx+0.1, 0.67, 3.3, 0.28, title,
                 font_size=11, bold=True, color=fg)
        for j, item in enumerate(items):
            ty = 1.0 + j * 1.44
            add_rect(slide, lx+0.1, ty, 3.3, 1.35, C_WHITE)
            add_text(slide, lx+0.2, ty+0.07, 3.1, 1.20,
                     f"{'①②③'[j]}  {item}", font_size=10, color=C_TEXT)


# ============================================================
# P9: 改善提案
# ============================================================
def build_p9_proposals(slide, d):
    slide_header(slide, "改善提案（優先度順）")
    proposals = d.get("proposals",[])

    for i, p in enumerate(proposals):
        ck = p.get("color","blue")
        bg, fg, icon = PROPOSAL_COLORS.get(ck, (C_LIGHT, C_TEXT, "●"))
        ty = 0.62 + i * 1.47
        add_rect(slide, 0.2, ty, 7.05, 1.38, bg)
        add_rect(slide, 0.2, ty, 1.0, 1.38, fg)
        add_text(slide, 0.2, ty+0.44, 1.0, 0.42, icon,
                 font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 1.3, ty+0.06, 5.9, 0.35,
                 p.get("title",""), font_size=10, bold=True, color=fg)
        add_text(slide, 1.3, ty+0.45, 5.9, 0.88,
                 p.get("body",""), font_size=10, color=C_TEXT)


# ============================================================
# P10: ページ別指標（2ヶ月並列）
# ============================================================
def build_p10_pages(slide, d):
    slide_header(slide, "ページ別指標（GA4）2ヶ月並列")
    for idx, key in enumerate(["page_metrics_1st","page_metrics_2nd"]):
        pages = d.get(key,[])
        if not pages: continue
        lx = 0.2 + idx * 3.6
        label = pages[0].get("ym","") if pages else ""
        add_text(slide, lx, 0.62, 3.45, 0.25, f"■ {label}",
                 font_size=9, bold=True, color=C_PRIMARY)
        h = ["ページパス","PV","滞在時間","ユーザー"]
        r = [[p["page_path"], f"{p['pageviews']:,}", p["duration"],
              f"{p['total_users']:,}"] for p in pages]
        add_table(slide, lx, 0.90, 3.45, 4.10, h, r[:10], fs=8,
                  col_widths=[1.55, 0.5, 0.8, 0.6])


# ============================================================
# P11-P12: 流入元
# ============================================================
def build_traffic_slide(slide, d, key, month_label):
    slide_header(slide, f"分析：流入（{month_label}）（GA4）")
    sources = d.get(key,[])
    if not sources: return
    label = sources[0].get("ym","")
    add_text(slide, 0.35, 0.62, 5.0, 0.25, f"■ {label}",
             font_size=9, bold=True, color=C_PRIMARY)
    h = ["参照元/メディア","セッション","セッション前月差分","ユーザー","ユーザー前月差分"]
    r = [[s["source_medium"],
          f"{s['sessions']:,}",
          delta_str(s.get("sessions_delta")),
          f"{s['total_users']:,}",
          delta_str(s.get("total_users_delta"))] for s in sources]
    # table width extended to 7.1
    add_table(slide, 0.2, 0.90, 7.1, 4.10, h, r[:10], fs=9,
              col_widths=[2.8, 1.0, 1.3, 1.0, 1.0])


# ============================================================
# P13-P14: 商圏内流入
# ============================================================
def build_area_slide(slide, d, key, month_label):
    slide_header(slide, f"分析：商圏内流入（{month_label}）（GA4）")
    sources = d.get(key,[])
    if not sources: return
    label = sources[0].get("ym", month_label)
    
    # Check if we have multiple cities by seeing if 'city' exists in rows
    # We display a generic target area title instead of a single parsed city.
    add_text(slide, 0.35, 0.62, 7.0, 0.25,
             f"■ {label}　対象エリア: 設定された各商圏",
             font_size=9, bold=True, color=C_PRIMARY)
    
    h = ["セッションの参照元/メディア", "対象エリア", "セッション", "セッション前月差分", "ユーザー", "ユーザー前月差分"]
    r = []
    for s in sources:
        r.append([
            s["source_medium"],
            s.get("city", "-"),
            f"{s['sessions']:,}",
            delta_str(s.get("sessions_delta")),
            f"{s['total_users']:,}",
            delta_str(s.get("total_users_delta"))
        ])
            
    add_table(slide, 0.2, 0.90, 7.1, 4.10, h, r[:15], fs=8,
              col_widths=[2.0, 1.5, 0.9, 0.9, 0.9, 0.9])


# ============================================================
def make_chart_transparent(chart):
    # For chart shape (graphic frame)
    spPr = chart._chartSpace.find(qn('c:spPr'))
    if spPr is None:
        spPr = chart._chartSpace.makeelement(qn('c:spPr'))
        chart._chartSpace.append(spPr)
    noFill = spPr.find(qn('a:noFill'))
    if noFill is None:
        noFill = spPr.makeelement(qn('a:noFill'))
        spPr.append(noFill)
        
    # For PlotArea
    plotArea = chart._chartSpace.find(qn('c:plotArea'))
    if plotArea is not None:
        p_spPr = plotArea.find(qn('c:spPr'))
        if p_spPr is None:
            p_spPr = plotArea.makeelement(qn('c:spPr'))
            plotArea.insert(0, p_spPr)
        if p_noFill is None:
            p_noFill = p_spPr.makeelement(qn('a:noFill'))
            p_spPr.append(p_noFill)

def make_overlay_invisible(chart):
    from pptx.oxml.ns import qn
    chart.has_title = True
    if chart.chart_title:
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = ""
    chart.value_axis.visible = False
    chart.category_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False
    for ax in [chart.value_axis._element, chart.category_axis._element]:
        tick = ax.find(qn('c:tickLblPos'))
        if tick is not None: tick.set('val', 'none')
        sp = ax.find(qn('c:spPr'))
        if sp is not None: ax.remove(sp)
        tx = ax.find(qn('c:txPr'))
        if tx is not None: ax.remove(tx)

def format_chart_axes(chart, fs=9):
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = ""
    
    if chart.has_legend:
        chart.legend.font.size = Pt(fs)
        chart.legend.font.name = FONT_JP
    if hasattr(chart, 'value_axis') and chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(fs)
        chart.value_axis.tick_labels.font.name = FONT_JP
    if hasattr(chart, 'category_axis') and chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(fs)
        chart.category_axis.tick_labels.font.name = FONT_JP
    for s in chart.series:
        if hasattr(s, 'has_data_labels') and s.has_data_labels:
            try:
                s.data_labels.font.size = Pt(fs)
                s.data_labels.font.name = FONT_JP
            except Exception:
                pass

# ============================================================
# P15: 広告基本指標 (Ads) 月次
# ============================================================
def build_p15_ads_monthly(slide, d):
    slide_header(slide, "広告基本指標 (Ads)")
    ads_m = d.get("ads_monthly", [])
    if not ads_m:
        add_text(slide, 0.35, 1.0, 5.0, 0.5, "データがありません", font_size=12, color=C_TEXT)
        return

    # Sort so oldest is first for charts, or newest first for table?
    # Reference image shows oldest first (2025-9月, 10, 11, 12, 1) in table and charts.
    ads_m = sorted(ads_m, key=lambda x: x["ym_raw"])

    # --- テーブル ---
    h = ["年月（日付）", "費用", "CV", "CPA", "クリック", "CPC", "CTR", "表示回数", "CVR"]
    r = []
    for m in ads_m:
        r.append([
            m["ym"],
            f"{m.get('cost',0):,.2f}" if isinstance(m.get('cost'), float) else f"{m.get('cost',0):,}",
            f"{m.get('cv',0):,.2f}",
            f"{m.get('cpa',0):,.0f}" if m.get('cv',0) > 0 else "",
            f"{m.get('clicks',0):,}",
            f"{m.get('cpc',0):,.0f}",
            f"{m.get('ctr',0):.2f}%" if m.get('ctr',0) > 2 else f"{m.get('ctr',0)*100:.2f}%" if m.get('ctr',0) < 1 else f"{m.get('ctr',0):.2f}%", 
            f"{m.get('impressions',0):,}",
            f"{m.get('cvr',0):.2f}%" if m.get('cvr',0) > 2 else f"{m.get('cvr',0)*100:.2f}%" if m.get('cvr',0) < 1 else f"{m.get('cvr',0):.2f}%"
        ])
    add_table(slide, 0.2, 0.6, 7.1, 1.0, h, r, fs=8,
              col_widths=[1.1, 0.9, 0.6, 0.8, 0.8, 0.6, 0.7, 0.9, 0.7])

    categories = [m["ym"].replace("年", "/").replace("月", "") for m in ads_m]

    # --- 左側グラフ: コスト (Bar) ---
    add_text(slide, 0.2, 2.0, 3.0, 0.3, "コスト (Ads)\nby Month", font_size=9, color=C_SUBTEXT)
    fpath1 = "temp_chart_p15_cost.png"
    plot_utils.save_bar_chart(
        categories, [m.get("cost", 0) for m in ads_m], "Cost",
        fpath1, width=320, height=220
    )
    slide.shapes.add_picture(fpath1, inch(0.2), inch(2.4), width=inch(3.3))

    # --- 右側グラフ: CV・CVR (Line+Bar combo) ---
    add_text(slide, 3.6, 2.0, 3.0, 0.3, "CV・CVR (Ads)\nby Month", font_size=9, color=C_SUBTEXT)
    fpath2 = "temp_chart_p15_cv_cvr.png"
    plot_utils.save_combo_chart(
        categories,
        [m["cv"] for m in ads_m], "CV",
        [m["cvr"] for m in ads_m], "CVR",
        fpath2, width=320, height=220
    )
    slide.shapes.add_picture(fpath2, inch(3.6), inch(2.4), width=inch(3.3), height=inch(2.8)) # Light blue


# ============================================================
# P16: 広告基本指標_週次 (Ads)
# ============================================================
def build_p16_ads_weekly(slide, d):
    slide_header(slide, "広告基本指標_週次 (Ads)")
    ads_w = d.get("ads_weekly", [])
    if not ads_w:
        add_text(slide, 0.35, 1.0, 5.0, 0.5, "データがありません", font_size=12, color=C_TEXT)
        return

    ads_w = sorted(ads_w, key=lambda x: x["week"])
    
    h = ["週毎（日付）", "費用", "CV", "CPA", "クリック", "CPC", "CTR", "表示回数", "CVR"]
    r = []
    for m in ads_w:
        r.append([
            m["week"],
            f"{m.get('cost',0):,.2f}" if isinstance(m.get('cost'), float) else f"{m.get('cost',0):,}",
            f"{m.get('cv',0):,.2f}",
            f"{m.get('cpa',0):,.0f}" if m.get('cv',0) > 0 else "",
            f"{m.get('clicks',0):,}",
            f"{m.get('cpc',0):,.0f}",
            f"{m.get('ctr',0):.2f}%" if m.get('ctr',0) > 2 else f"{m.get('ctr',0)*100:.2f}%" if m.get('ctr',0) < 1 else f"{m.get('ctr',0):.2f}%",
            f"{m.get('impressions',0):,}",
            f"{m.get('cvr',0):.2f}%" if m.get('cvr',0) > 2 else f"{m.get('cvr',0)*100:.2f}%" if m.get('cvr',0) < 1 else f"{m.get('cvr',0):.2f}%"
        ])
    # Table heights might overflow if many weeks. Assume ~10 weeks.
    add_table(slide, 0.2, 0.5, 7.1, 1.0, h, r, fs=7,
              col_widths=[1.1, 0.8, 0.6, 0.6, 0.7, 0.6, 0.7, 1.1, 0.9])

    categories = [m["week"][-5:] + " " for m in ads_w] # use short dates, append space to make Plotly treat as string

    # --- 左側グラフ: 表示回数-クリック数 (Click=Bar, Imp=Line overlay) ---
    add_text(slide, 0.2, 2.8, 3.0, 0.2, "表示回数-クリック数 (Ads)", font_size=9, color=C_SUBTEXT)
    fpath1 = "temp_chart_p16_ct_imp.png"
    plot_utils.save_combo_chart(
        categories,
        [m["clicks"] for m in ads_w], "Click",
        [m["impressions"] for m in ads_w], "Imp",
        fpath1, width=320, height=180
    )
    slide.shapes.add_picture(fpath1, inch(0.2), inch(3.1), width=inch(3.3))

    # --- 右側グラフ: CV-CPA ---
    add_text(slide, 3.6, 2.8, 3.0, 0.2, "CV-CPA (Ads)", font_size=9, color=C_SUBTEXT)
    fpath2 = "temp_chart_p16_cpa_cv.png"
    plot_utils.save_combo_chart(
        categories,
        [m["cpa"] for m in ads_w], "CPA",
        [m["cv"] for m in ads_w], "CV",
        fpath2, width=320, height=180
    )
    slide.shapes.add_picture(fpath2, inch(3.6), inch(3.1), width=inch(3.3), height=inch(2.0))


# ============================================================
# P17: キャンペーン指標推移 (Ads)
# ============================================================
def build_p17_ads_campaign(slide, d):
    slide_header(slide, "キャンペーン指標推移（費用）")
    ads_c = d.get("ads_campaigns", [])
    if not ads_c:
        add_text(slide, 0.35, 1.0, 5.0, 0.5, "データがありません", font_size=12, color=C_TEXT)
        return

    # To plot line chart, we need costs grouped by campaign over YM
    cd = ChartData()
    
    # Get sorted distinct ym
    ym_set = sorted(list(set(m["ym_raw"] for m in ads_c)))
    cd.categories = ym_set
    
    # Group by campaign
    camp_dict = {}
    for m in ads_c:
        cname = m["campaign"]
        if cname not in camp_dict:
            camp_dict[cname] = {ym: 0 for ym in ym_set}
        camp_dict[cname][m["ym_raw"]] += float(m.get("cost", 0))
    
    # Add series for top N campaigns by total cost to avoid overly messy chart
    top_camps = sorted(list(camp_dict.keys()), key=lambda c: sum(camp_dict[c].values()), reverse=True)[:5]
    for c in top_camps:
        cd.add_series(c, [camp_dict[c][ym] for ym in ym_set])

    add_text(slide, 0.2, 0.45, 5.0, 0.25, "キャンペーン指標推移（費用）",
             font_size=10, bold=True, color=C_TEXT)

    chart_costs = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, inch(0.2), inch(0.7), inch(7.1), inch(2.4), cd
    ).chart
    chart_costs.has_legend = True
    chart_costs.legend.position = XL_LEGEND_POSITION.BOTTOM
    format_chart_axes(chart_costs, fs=8)
    chart_costs.plots[0].has_data_labels = False

    # --- 下部: キャンペーン別 最新月テーブル ---
    add_text(slide, 0.2, 3.3, 4.0, 0.25, "広告基本指標_週次・キャンペーン (Ads)", font_size=9, color=C_SUBTEXT)
    h = ["年月", "キャンペーン名", "費用", "CV", "CPA", "クリック", "CPC", "CTR", "表示回数", "CVR"]
    r = []
    
    # Filter for the latest month in data
    latest_ym = max(ym_set) if ym_set else ""
    latest_camps = [m for m in ads_c if m["ym_raw"] == latest_ym]
    
    for m in latest_camps:
        short_ym = m["ym"].replace("年", "/").replace("月", "")
        r.append([
            short_ym,
            m["campaign"],
            f"{m.get('cost',0):,.0f}",
            f"{m.get('cv',0):.2f}",
            f"{m.get('cpa',0):,.0f}" if m.get('cpa') else "",
            f"{m.get('clicks',0):,}",
            f"{m.get('cpc',0):,.0f}" if m.get('cpc') else "",
            f"{m.get('ctr',0)*100:.2f}%",
            f"{m.get('impressions',0):,}",
            f"{m.get('cvr',0):.2f}%" if m.get('cvr',0) > 2 else f"{m.get('cvr',0)*100:.2f}%" if m.get('cvr',0) < 1 else f"{m.get('cvr',0):.2f}%"
        ])
    add_table(slide, 0.2, 3.5, 7.1, 1.5, h, r, fs=7,
              col_widths=[0.6, 2.0, 0.7, 0.45, 0.65, 0.5, 0.45, 0.55, 0.7, 0.5])


# メイン
# ============================================================
def main():
    input_path  = sys.argv[1] if len(sys.argv) > 1 else "report_data_test.json"
    output_path = sys.argv[2] if len(sys.argv) > 2 else "report_output.pptx"

    print(f"📂 読み込み: {input_path}")
    d = load_data(input_path)

    print(f"📄 テンプレート読み込み: {TEMPLATE_FILE}")
    prs = Presentation(TEMPLATE_FILE)

    # P1,P2はテンプレートを一切変更しない（固定）
    print("P1,P2 固定（変更なし）")

    # テンプレートSlide3を参照用に保持
    tmpl = prs.slides[2]

    def new_slide():
        return add_content_slide(prs, tmpl)

    print("P3 CV・CPA 生成中...")
    build_p3_cv(new_slide(), d)

    print("P4 全体サマリー 生成中...")
    build_p4_summary(new_slide(), d)

    print("P5 詳細サマリー 生成中...")
    build_p5_detail(new_slide(), d)

    print("P6 月別GA4分析 生成中...")
    build_p6_ga4(new_slide(), d)

    print("P7 月別GSC分析 生成中...")
    build_p7_gsc(new_slide(), d)

    print("P7.5 商圏GSC分析 生成中...")
    build_p7_5_gsc_area(new_slide(), d, new_slide_fn=new_slide)

    print("P8 当月分析 生成中...")
    build_p8_analysis(new_slide(), d)

    print("P9 改善提案 生成中...")
    build_p9_proposals(new_slide(), d)

    print("P10 ページ別指標 生成中...")
    build_p10_pages(new_slide(), d)

    print("P11 流入_当月 生成中...")
    build_traffic_slide(new_slide(), d, "traffic_sources_1st", "当月")

    print("P12 流入_前月 生成中...")
    build_traffic_slide(new_slide(), d, "traffic_sources_2nd", "前月")

    print("P13 商圏内流入_当月 生成中...")
    build_area_slide(new_slide(), d, "area_traffic_1st", "当月")

    print("P14 商圏内流入_前月 生成中...")
    build_area_slide(new_slide(), d, "area_traffic_2nd", "前月")

    if d.get("ads_monthly"):
        print("P15 広告基本指標(月次) 生成中...")
        build_p15_ads_monthly(new_slide(), d)

        print("P16 広告基本指標(週次) 生成中...")
        build_p16_ads_weekly(new_slide(), d)

        print("P17 キャンペーン指標推移 生成中...")
        build_p17_ads_campaign(new_slide(), d)

    # テンプレートのSlide3（雛形）を最後に削除
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    # Slide3はインデックス2（0始まり）→ 今は14枚追加後は位置2にある
    # 追加後はslide3が先頭側にあるので削除
    slide3_rId = sldIdLst[2].get("r:id")
    slide3_elem = sldIdLst[2]
    sldIdLst.remove(slide3_elem)
    # 参照パーツも削除
    try:
        del prs.slides.part.related_parts[slide3_rId]
    except: pass

    prs.save(output_path)
    print(f"\n✅ 完成! → {output_path}")
    print(f"   スライド数: {len(prs.slides)}枚")

if __name__ == "__main__":
    main()

def generate(data: dict, template_path: str, output_path: str):
    """APIから呼び出し可能なPPTX生成エントリーポイント"""
    import os, tempfile, json

    # dataをtmpに書き出してload_dataで読み込む
    tmp = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
    json.dump(data, tmp, ensure_ascii=False)
    tmp.close()

    try:
        d = load_data(tmp.name)
    finally:
        os.unlink(tmp.name)

    global TEMPLATE_FILE
    TEMPLATE_FILE = template_path
    prs = Presentation(template_path)

    tmpl = prs.slides[2]

    def new_slide():
        return add_content_slide(prs, tmpl)

    build_p3_cv(new_slide(), d)
    build_p4_summary(new_slide(), d)
    build_p5_detail(new_slide(), d)
    build_p6_ga4(new_slide(), d)
    build_p7_gsc(new_slide(), d)
    build_p7_5_gsc_area(new_slide(), d, new_slide_fn=new_slide)
    build_p8_analysis(new_slide(), d)
    build_p9_proposals(new_slide(), d)
    build_p10_pages(new_slide(), d)
    build_traffic_slide(new_slide(), d, "traffic_sources_1st", "当月")
    build_traffic_slide(new_slide(), d, "traffic_sources_2nd", "前月")
    build_area_slide(new_slide(), d, "area_traffic_1st", "当月")
    build_area_slide(new_slide(), d, "area_traffic_2nd", "前月")
    if d.get("ads_monthly"):
        build_p15_ads_monthly(new_slide(), d)
        build_p16_ads_weekly(new_slide(), d)
        build_p17_ads_campaign(new_slide(), d)

    # テンプレートのSlide3（雛形）を削除
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    slide3_rId = sldIdLst[2].get("r:id")
    slide3_elem = sldIdLst[2]
    sldIdLst.remove(slide3_elem)
    try:
        del prs.slides.part.related_parts[slide3_rId]
    except: pass

    prs.save(output_path)
